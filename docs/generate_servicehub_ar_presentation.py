from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.dml import MSO_THEME_COLOR

OUT = '/home/server/projects/servicehub/docs/ServiceHub-Arabic-Presentation.pptx'
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(12, 20, 20)
SURFACE = RGBColor(22, 38, 38)
ACCENT = RGBColor(88, 180, 150)
GOLD = RGBColor(230, 180, 90)
TEXT = RGBColor(242, 246, 244)
MUTED = RGBColor(177, 196, 190)
RED = RGBColor(232, 116, 116)
FONT = 'DejaVu Sans'


def rtl(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set('rtl', '1')


def box(slide, x, y, w, h, text='', size=18, color=TEXT, bold=False, align=PP_ALIGN.RIGHT, fill=None, line=None, radius=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill if fill else BG
    shape.line.color.rgb = line if line else (fill if fill else BG)
    tf = shape.text_frame
    tf.clear(); tf.margin_left = Inches(.18); tf.margin_right = Inches(.18); tf.margin_top = Inches(.10); tf.margin_bottom = Inches(.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text; p.alignment = align; rtl(p)
    for run in p.runs:
        run.font.name = FONT; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return shape


def text(slide, x, y, w, h, value, size=18, color=TEXT, bold=False, align=PP_ALIGN.RIGHT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True; tf.margin_left = 0; tf.margin_right = 0
    p = tf.paragraphs[0]; p.text = value; p.alignment = align; rtl(p)
    for run in p.runs:
        run.font.name = FONT; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return tb


def bullet_list(slide, x, y, w, h, items, size=17):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True; tf.margin_left = Inches(.05); tf.margin_right = Inches(.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = '• ' + item; p.alignment = PP_ALIGN.RIGHT; rtl(p); p.space_after = Pt(9)
        for r in p.runs:
            r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = TEXT
    return tb


def base(title, subtitle=None, n=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = BG
    box(slide, 0, 0, 13.333, .12, fill=ACCENT, radius=False)
    text(slide, .65, .42, 12.0, .45, title, 27, TEXT, True)
    if subtitle: text(slide, .65, .92, 12.0, .32, subtitle, 12, MUTED)
    if n is not None: text(slide, .5, 7.08, 1.0, .2, str(n).zfill(2), 10, MUTED, align=PP_ALIGN.LEFT)
    text(slide, 9.8, 7.05, 2.9, .22, 'ServiceHub • عرض المنتج', 10, MUTED)
    return slide

# 1 cover
s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
box(s, 0, 0, 13.333, .14, fill=ACCENT, radius=False)
text(s, .8, 1.35, 11.8, .75, 'ServiceHub', 48, TEXT, True, PP_ALIGN.CENTER)
text(s, 1.2, 2.25, 10.9, .6, 'منصة موحّدة لاكتشاف الخدمات، الحجز، الدفع، وإدارة الأعمال', 25, ACCENT, True, PP_ALIGN.CENTER)
text(s, 1.4, 3.12, 10.5, .55, 'عرض عربي مفصل للنظام، الصفحات، المزايا، وخارطة التطوير', 18, MUTED, align=PP_ALIGN.CENTER)
for i,(label,color) in enumerate([('العميل',ACCENT),('مزود الخدمة',GOLD),('الإدارة',RGBColor(142,167,255))]):
    box(s, 2.1+i*3.15, 4.65, 2.55, .75, label, 19, BG, True, PP_ALIGN.CENTER, color, color)
text(s, .8, 6.6, 11.7, .35, 'نسخة عرض وتجربة داخلية • يوليو 2026', 13, MUTED, align=PP_ALIGN.CENTER)

# 2 problem
s=base('المشكلة التي يحلها ServiceHub','تجربة الحجز في الخدمات المحلية غالباً موزعة وغير قابلة للقياس.',2)
bullet_list(s,.8,1.55,5.8,4.8,['العميل يتنقل بين رسائل واتساب ومكالمات ودفاتر مواعيد.', 'مزود الخدمة لا يرى صورة واحدة للحجوزات والخدمات والتوفر.', 'الإدارة تحتاج موافقات، مؤشرات تشغيل، ونزاعات قابلة للمتابعة.', 'الحجز المتزامن قد ينتج عنه تعارضات مكلفة إن لم يكن محمياً على مستوى قاعدة البيانات.'],19)
box(s,7.15,1.55,5.2,4.45,'النتيجة المطلوبة\n\nسوق خدمات متعدد المزودين يربط الاكتشاف بالحجز والدفع والمراجعة ضمن رحلة واحدة واضحة.',23,TEXT,True,PP_ALIGN.CENTER,SURFACE,ACCENT)

# 3 value
s=base('الفكرة والقيمة','منصة متعددة المزودين لقطاعات مثل المطاعم والصالونات والاستشارات والصيانة.',3)
for x,t,d,c in [(8.65,'اكتشف','بحث وفلاتر وخدمات ومراجعات',ACCENT),(5.38,'احجز','مواعيد وتوفر وحجز آمن',GOLD),(2.1,'أدر','لوحات مزود وإدارة وتشغيل',RGBColor(142,167,255))]:
    box(s,x,1.65,2.55,1.2,t,22,BG,True,PP_ALIGN.CENTER,c,c); text(s,x,3.05,2.55,.7,d,14,TEXT,align=PP_ALIGN.CENTER)
text(s,.9,4.65,11.6,.65,'القيمة ليست في عرض قائمة خدمات فقط؛ بل في تحويل العملية كاملة إلى بيانات قابلة للتشغيل والقياس.',20,TEXT,True,PP_ALIGN.CENTER)

# 4 roles
s=base('الأدوار الثلاثة في النظام','كل دور يملك صلاحيات ومسارات مختلفة، مع فصل واضح للبيانات.',4)
for i,(h,items,c) in enumerate([('العميل',['اكتشاف ومقارنة','اختيار موعد','دفع وتقييم'],ACCENT),('مزود الخدمة',['إدارة الخدمات','تحديد التوفر','متابعة الحجوزات'],GOLD),('الإدارة',['اعتماد المزودين','مؤشرات وتقارير','متابعة النزاعات'],RGBColor(142,167,255))]):
    x=8.85-i*3.85; box(s,x,1.45,3.3,4.55,'',fill=SURFACE,line=c); text(s,x+.2,1.75,2.9,.45,h,22,c,True,PP_ALIGN.CENTER); bullet_list(s,x+.25,2.55,2.8,2.5,items,16)

# 5 journey
s=base('رحلة العميل: من البحث إلى التقييم','مسار واضح يقلل الاحتكاك ويوثق كل نقطة قرار.',5)
steps=['اكتشاف','صفحة المزود','اختيار الموعد','تأكيد الدفع','الحجز','التقييم']
for i,step in enumerate(steps):
    x=.75+i*2.08; box(s,x,2.2,1.65,.7,step,14,BG,True,PP_ALIGN.CENTER,ACCENT,ACCENT)
    if i<5: text(s,x+1.63,2.34,.35,.25,'←',18,ACCENT,True,PP_ALIGN.CENTER)
bullet_list(s,1.25,3.6,10.9,2.1,['بحث حسب التصنيف، الاسم، السعر، والتقييم.', 'حجز محدد الوقت مع hold مدته 5 دقائق أثناء الدفع.', 'إلغاء منضبط بسياسة زمنية، ثم مراجعة من 1 إلى 5 نجوم بعد الحجز المؤكد.'],18)

# 6 pages customer
s=base('صفحات العميل','الواجهة الأمامية تغطي رحلة العميل الأساسية من البداية للنهاية.',6)
items=[('الرئيسية','اكتشاف وفلاتر ومزودون'),('صفحة المزود','الخدمات، التفاصيل، والمراجعات'),('الحجز','تاريخ + شبكة المواعيد المتاحة'),('الدفع','تأكيد محاكاة الدفع أو Stripe'),('حجوزاتي','الحالة، الإلغاء، والتقييم'),('الإشعارات','تنبيهات الحجز والدفع')]
for i,(a,b) in enumerate(items):
    x=8.7-(i%3)*3.9; y=1.45+(i//3)*2.2; box(s,x,y,3.25,1.45,'',fill=SURFACE,line=ACCENT); text(s,x+.2,y+.25,2.85,.34,a,18,ACCENT,True,PP_ALIGN.CENTER); text(s,x+.22,y+.76,2.8,.34,b,13,TEXT,align=PP_ALIGN.CENTER)

# 7 vendor pages
s=base('صفحات مزود الخدمة','لوحة تشغيل تستبدل الجداول المتناثرة بمصدر بيانات واحد.',7)
bullet_list(s,.75,1.45,5.8,4.9,['لوحة KPI: حجوزات اليوم، الإيراد المؤكد، الإلغاءات، والخدمات الأكثر طلباً.', 'الخدمات: إضافة وتعديل وتعطيل الخدمة مع السعر والمدة والتصنيف.', 'الجدول الأسبوعي: تحديد فترات التوفر وحفظها في API، لا في localStorage.', 'الحجوزات: عرض حجوزات المزود وإدارة الإلغاء ضمن الصلاحيات.'],18)
box(s,7.1,1.45,5.3,4.6,'تحسين حديث\n\nتم ربط لوحة Vendor بمصادر API حقيقية، مع حالات تحميل وخطأ وفراغ، وإزالة بيانات العرض المؤقتة.',21,TEXT,True,PP_ALIGN.CENTER,SURFACE,GOLD)

# 8 admin
s=base('صفحة الإدارة','طبقة تشغيل وإشراف وليست مجرد شاشة تقارير.',8)
for i,(a,b) in enumerate([('المزودون','اعتماد، تعليق، ومتابعة الحالة'),('المؤشرات','المستخدمون، GMV، والعمولة'),('التقارير','اتجاه الإيرادات وأفضل المزودين'),('النزاعات','قائمة الحجوزات الملغاة وحالاتها')]):
    x=8.6-(i%2)*4.25; y=1.45+(i//2)*2.2; box(s,x,y,3.65,1.45,'',fill=SURFACE,line=RGBColor(142,167,255)); text(s,x+.2,y+.25,3.25,.35,a,18,RGBColor(142,167,255),True,PP_ALIGN.CENTER); text(s,x+.2,y+.78,3.25,.35,b,13,TEXT,align=PP_ALIGN.CENTER)

# 9 architecture
s=base('المعمارية التقنية','فصل واضح بين الواجهة، API، وقاعدة البيانات مع وحدات قابلة للاختبار.',9)
for x,h,d,c in [(9.1,'Next.js 14','واجهة App Router + Tailwind',ACCENT),(5.0,'NestJS 11','REST API + وحدات المجال',GOLD),(.9,'PostgreSQL 16','Prisma + migrations + constraints',RGBColor(142,167,255))]:
    box(s,x,2.15,3.25,1.25,h,21,BG,True,PP_ALIGN.CENTER,c,c); text(s,x,3.55,3.25,.45,d,13,TEXT,align=PP_ALIGN.CENTER)
text(s,4.05,2.48,.75,.35,'←',25,MUTED,True,PP_ALIGN.CENTER); text(s,8.1,2.48,.75,.35,'←',25,MUTED,True,PP_ALIGN.CENTER)
bullet_list(s,1.1,4.65,11.1,1.5,['Auth: JWT في HttpOnly cookies + CSRF للطلبات المتغيرة.', 'Payments: provider abstraction بين Stripe للإنتاج وMock للتجربة.', 'Notifications: REST polling موثق كاختيار MVP.'],16)

# 10 safety
s=base('موثوقية الحجز والأمان','نقاط الحماية التي تجعل النظام مناسباً للتجربة الجادة لا للعرض الشكلي فقط.',10)
items=[('منع double booking','PostgreSQL EXCLUDE constraint هو مصدر الحقيقة حتى تحت التزامن.'),('صلاحيات وعزل بيانات','Customer / Vendor / Admin مع tenant-scoped reads وIDOR checks.'),('جلسات آمنة','HttpOnly access/refresh cookies؛ لا JWT قابل للقراءة في JavaScript.'),('CSRF','double-submit token للطلبات POST/PUT/PATCH/DELETE في cookie-auth.')]
for i,(a,b) in enumerate(items):
    y=1.35+i*1.25; box(s,7.75,y,4.55,.88,a,17,BG,True,PP_ALIGN.CENTER,ACCENT,ACCENT); text(s,.95,y+.12,6.35,.5,b,14,TEXT)

# 11 demo decision
s=base('حدود العرض الحالية بوضوح','الشفافية جزء من جودة المنتج.',11)
box(s,.85,1.45,11.65,1.15,'Google (Demo) محاكاة محلية فقط — لا OAuth حقيقي، لا بيانات اعتماد Google، ولا هوية خارجية.',19,TEXT,True,PP_ALIGN.CENTER,SURFACE,GOLD)
bullet_list(s,1.1,3.05,11.1,2.6,['الدفع يمكن تجربته عبر Mock provider؛ Stripe هو مسار الإنتاج عند إضافة الأسرار وwebhook.', 'التنبيهات تعمل عبر REST polling، وليست WebSocket حالياً.', 'الروابط الحالية للاختبار داخل tailnet؛ ليست نشر إنترنت عام.'],17)

# 12 evidence
s=base('ما تم التحقق منه حتى الآن','الأرقام هنا ناتجة عن تنفيذ واختبارات فعلية في نسخة العرض.',12)
for i,(n,label,c) in enumerate([('91/91','اختبارات API ناجحة',ACCENT),('16','صفحة Web تُبنى بنجاح',GOLD),('3','أدوار أساسية',RGBColor(142,167,255))]):
    x=8.8-i*3.85; box(s,x,1.6,3.2,1.7,'',fill=SURFACE,line=c); text(s,x+.15,1.88,2.9,.48,n,28,c,True,PP_ALIGN.CENTER); text(s,x+.15,2.55,2.9,.3,label,13,TEXT,align=PP_ALIGN.CENTER)
bullet_list(s,1.1,4.4,11.1,1.7,['API وWeb production builds ناجحة.', 'Vendor dashboard تحقق من tenant isolation، بما فيه اختبار بيانات Vendor ثانٍ.', 'اختبارات CSRF تثبت الرفض بدون token والنجاح مع token صحيح.'],16)

# 13 testing plan
s=base('خطة الاختبار القادمة','بعد تجربة المنتج، نتحول من بناء المزايا إلى evidence قبل أي ادعاء Release جاهز.',13)
for i,(a,b) in enumerate([('1. تجربة يدوية','تجربة العميل، Vendor، Admin وتسجيل الملاحظات.'),('2. قبول API','صلاحيات، transitions، validation، أخطاء متوقعة.'),('3. E2E','رحلات الحجز والدفع والإلغاء من المتصفح.'),('4. أداء وRace','التزامن، زمن الاستجابة، وقياس baseline.')]):
    y=1.25+i*1.25; text(s,8.35,y,3.6,.35,a,18,ACCENT,True); text(s,.9,y,7.0,.45,b,16,TEXT)

# 14 roadmap
s=base('المزايا المستقبلية','مسار واضح بعد تثبيت النسخة الأساسية واختبارها.',14)
items=[('i18n كامل','Arabic RTL وEnglish LTR مع language switcher ومسارات locale.'),('Lifecycle وتنبيهات','Reminders 24h/1h، dedup، وإثبات booking race E2E.'),('تجاري','إعداد العمولة، CSV، طلبات payout، ومراجعة البلاغات.'),('مراسلة','رسائل مرتبطة بالحجز مع صلاحيات participants وpagination.'),('تشغيل','Swagger، CI، Docker، backups، metrics وrelease acceptance.')]
for i,(a,b) in enumerate(items):
    y=1.2+i*1.03; box(s,9.3,y,2.55,.65,a,15,BG,True,PP_ALIGN.CENTER,GOLD,GOLD); text(s,.9,y+.1,7.9,.35,b,14,TEXT)

# 15 demo script
s=base('سيناريو التجربة المقترح','ابدأ بالمسار الذي يظهر قيمة المنصة بأسرع شكل.',15)
bullet_list(s,.95,1.35,11.3,4.9,['افتح الرابط وسجل عبر Google (Demo) — ستظهر عبارة توضح أنها محاكاة فقط.', 'اكتشف مزوداً، افتح صفحة خدماته، واختر موعداً متاحاً.', 'أنشئ حجزاً ثم جرّب Mock payment confirmation.', 'افتح حجوزاتي ثم الإشعارات؛ جرّب الإلغاء إن كان الموعد أبعد من 24 ساعة.', 'لـ Vendor/Admin: استخدم حساباً مناسباً في قاعدة التجربة عند توفره، ثم راقب dashboard والخدمات والجدول.'],18)

# 16 closing
s=base('الخلاصة','ServiceHub أصبح MVP متعدد الأدوار قابلاً للتجربة، مع مسار واضح لتحويله إلى Release Candidate مختبر.',16)
box(s,1.15,1.55,11.05,1.25,'الآن: تجربة المستخدم وجمع الملاحظات.\nبعدها: نثبت النطاق ونبني acceptance evidence قبل الإطلاق.',24,TEXT,True,PP_ALIGN.CENTER,SURFACE,ACCENT)
text(s,1.3,4.2,10.7,.4,'شكراً',30,ACCENT,True,PP_ALIGN.CENTER)
text(s,1.3,4.85,10.7,.35,'ServiceHub — خدمة أوضح، حجز أضمن، وتشغيل أذكى.',17,MUTED,align=PP_ALIGN.CENTER)

prs.save(OUT)
print(OUT)
