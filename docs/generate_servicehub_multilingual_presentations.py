"""Generate editable ServiceHub presentation decks in English, Arabic, and pt-BR.
All claims are grounded in repository documentation and the 2026-07-16 API test run.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

OUT_DIR = Path(__file__).parent
FONT = "DejaVu Sans"
BG = RGBColor(12, 24, 25)
SURFACE = RGBColor(21, 43, 43)
MINT = RGBColor(98, 205, 166)
GOLD = RGBColor(234, 187, 94)
LILAC = RGBColor(168, 178, 255)
TEXT = RGBColor(244, 248, 246)
MUTED = RGBColor(181, 201, 195)

DECKS = {
    "en": {
        "file": "ServiceHub-Presentation-EN.pptx", "rtl": False, "label": "PRODUCT PRESENTATION",
        "slides": [
            ("ServiceHub", "A multi-vendor services marketplace\nDiscovery · safe booking · payments · operations", ["Portfolio MVP · July 2026", "Customer · Vendor · Admin"]),
            ("The problem", "Local booking is fragmented", ["Customers move between messages, calls, and disconnected calendars.", "Providers lack one trusted operational view.", "Overlapping slots create expensive service failures.", "Administrators need approvals, metrics, and dispute visibility."]),
            ("The value proposition", "One connected service journey", ["Discover providers by category, price, rating, and service.", "Choose an available slot and create a protected timed hold.", "Complete payment through a provider boundary.", "Keep bookings, messaging, reviews, and updates in one place."]),
            ("Three product roles", "A focused experience for each participant", ["Customer: discovery, booking, cancellation, review, messages, notifications.", "Vendor: services, weekly availability, exceptions, reservations, operations.", "Admin: vendor status, marketplace KPIs, reports, disputes, and categories."]),
            ("Customer journey", "Designed as a clear, auditable flow", ["Search and inspect a provider profile.", "Pick a service, date, and available time slot.", "Reserve for five minutes while payment is in progress.", "Confirm, receive notifications, then review after an eligible booking."]),
            ("Booking reliability", "The database is the final authority", ["PostgreSQL EXCLUDE USING gist prevents overlapping active bookings per provider.", "The API also translates conflicts into a predictable client response.", "Expired payment holds are released by the domain lifecycle.", "The model is designed to remain correct under concurrent requests."]),
            ("Architecture", "Separated layers, testable domain modules", ["Next.js 14 web application for all role-based journeys.", "NestJS 11 REST API for auth, catalog, availability, booking, payment, and admin domains.", "Prisma 5 and PostgreSQL 16 for transactions, schema, and migrations.", "Mock provider for development; Stripe boundary for a properly configured production path."]),
            ("Security baseline", "Protection is part of the product contract", ["Role guards and ownership / IDOR checks protect sensitive actions.", "bcrypt password hashing, refresh-token revocation, throttling, and validation are in the API.", "CSRF guard coverage exists for cookie-authenticated mutation requests.", "Payment webhooks capture raw body and use idempotent state transitions."]),
            ("Quality evidence", "What was independently verified", ["16 July 2026: 15 API test suites passed.", "220 API tests passed with no failing snapshots.", "The evidence demonstrates tested behavior—not a public production deployment.", "Browser acceptance, external provider tests, and deployment evidence remain separate gates."]),
            ("Demo boundaries", "Clear claims build trust", ["PAYMENTS_PROVIDER=mock is for development and portfolio demonstrations only.", "Google (Demo) is a local simulation; it does not call Google and is not OAuth.", "Live Stripe requires secrets, verified webhooks, deployment configuration, and E2E testing.", "The current project is not claimed as a public marketplace handling real payments."]),
            ("Next evidence gates", "From delivered MVP to release candidate", ["Run browser acceptance in Arabic RTL and English LTR.", "Verify live Stripe test-mode webhooks and email delivery.", "Add CI, repeatable deploy, backup/restore, monitoring, and performance baselines.", "Collect user feedback before extending commercial scope."]),
            ("Thank you", "ServiceHub", ["Clearer service discovery. Safer booking. Smarter operations.", "Documentation is available in English, Arabic, and Portuguese (Brazil)."]),
        ],
    },
    "ar": {
        "file": "ServiceHub-Presentation-AR.pptx", "rtl": True, "label": "عرض المنتج",
        "slides": [
            ("ServiceHub", "منصة متعددة المزودين للخدمات\nاكتشاف · حجز آمن · دفع · تشغيل", ["MVP للـ portfolio · يوليو 2026", "عميل · مزود خدمة · إدارة"]),
            ("المشكلة", "الحجز المحلي غالباً عملية متفرقة", ["العميل ينتقل بين الرسائل والمكالمات والتقاويم المنفصلة.", "مزود الخدمة لا يملك رؤية تشغيلية موحدة.", "تعارض المواعيد قد يسبب فشلًا مكلفًا في الخدمة.", "الإدارة تحتاج اعتماد المزودين ومؤشرات ونزاعات واضحة."]),
            ("القيمة", "رحلة خدمة واحدة مترابطة", ["اكتشاف المزودين حسب الفئة والسعر والتقييم والخدمة.", "اختيار موعد متاح مع hold زمني محمي.", "الدفع عبر طبقة provider قابلة للتبديل.", "الحجوزات والرسائل والمراجعات والتحديثات في مكان واحد."]),
            ("الأدوار الثلاثة", "تجربة مركزة لكل طرف", ["العميل: اكتشاف، حجز، إلغاء، تقييم، رسائل، وإشعارات.", "المزود: خدمات، توفر أسبوعي، استثناءات، حجوزات، وتشغيل.", "الإدارة: حالة المزودين، مؤشرات السوق، تقارير، نزاعات، وفئات."]),
            ("رحلة العميل", "مسار واضح وقابل للمراجعة", ["البحث ثم قراءة صفحة مزود الخدمة.", "اختيار الخدمة والتاريخ والموعد المتاح.", "حجز مؤقت لخمس دقائق أثناء الدفع.", "تأكيد وإشعارات ثم تقييم بعد حجز مؤهل."]),
            ("موثوقية الحجز", "قاعدة البيانات هي مصدر الحقيقة", ["قيد PostgreSQL EXCLUDE USING gist يمنع تداخل الحجوزات النشطة للمزود.", "الـ API يحوّل التعارض إلى استجابة مفهومة للعميل.", "الحجوزات المؤقتة المنتهية تُحرّرها دورة حياة المجال.", "النموذج مصمم ليبقى صحيحًا مع الطلبات المتزامنة."]),
            ("المعمارية", "طبقات منفصلة ووحدات قابلة للاختبار", ["Next.js 14 للواجهات حسب الدور.", "NestJS 11 API لمجالات auth والخدمات والتوفر والحجز والدفع والإدارة.", "Prisma 5 وPostgreSQL 16 للمعاملات وschema وmigrations.", "Mock للتطوير وواجهة Stripe لمسار إنتاج مضبوط."]),
            ("خط الأساس الأمني", "الحماية جزء من عقد المنتج", ["Role guards وفحوص ownership وIDOR للمسارات الحساسة.", "bcrypt وrevocation وthrottling وvalidation داخل API.", "تغطية CsrfGuard لطلبات التعديل مع cookie-auth.", "Webhook raw body وحالات دفع idempotent."]),
            ("دليل الجودة", "ما تم التحقق منه فعليًا", ["16 يوليو 2026: نجاح 15 API test suites.", "نجاح 220 API tests دون snapshots فاشلة.", "هذا الدليل يثبت السلوك المختبَر، لا نشرًا عامًا.", "Browser QA والتكاملات الخارجية والنشر لها بوابات مستقلة."]),
            ("حدود العرض", "الوضوح يبني الثقة", ["PAYMENTS_PROVIDER=mock للتطوير والعرض فقط.", "Google (Demo) محاكاة محلية وليس OAuth حقيقيًا.", "Stripe الحقيقي يتطلب secrets وwebhooks ونشرًا وE2E.", "لا ادعاء بأن المنصة الآن سوق عام أو تعالج مدفوعات حقيقية."]),
            ("الخطوات التالية", "من MVP مكتمل إلى Release Candidate", ["Browser acceptance بالعربية RTL والإنجليزية LTR.", "اختبار Stripe test mode والبريد وwebhooks فعليًا.", "CI ونشر قابل للتكرار وbackup/restore وmonitoring وperformance.", "جمع feedback قبل توسيع النطاق التجاري."]),
            ("شكرًا", "ServiceHub", ["اكتشاف أوضح للخدمات. حجز أضمن. تشغيل أذكى.", "التوثيق متاح بالإنجليزية والعربية والبرتغالية البرازيلية."]),
        ],
    },
    "pt": {
        "file": "ServiceHub-Apresentacao-PT-BR.pptx", "rtl": False, "label": "APRESENTAÇÃO DO PRODUTO",
        "slides": [
            ("ServiceHub", "Marketplace de serviços com múltiplos fornecedores\nDescoberta · agendamento seguro · pagamentos · operação", ["MVP de portfólio · julho de 2026", "Cliente · Fornecedor · Administração"]),
            ("O problema", "Agendamento local é fragmentado", ["Clientes alternam entre mensagens, chamadas e calendários desconectados.", "Fornecedores não têm uma visão operacional confiável.", "Horários sobrepostos geram falhas caras no serviço.", "A administração precisa de aprovações, indicadores e disputas visíveis."]),
            ("A proposta", "Uma jornada de serviço conectada", ["Descoberta por categoria, preço, avaliação e serviço.", "Escolha de horário disponível com hold temporário protegido.", "Pagamento por uma camada de provider intercambiável.", "Reservas, mensagens, avaliações e atualizações no mesmo lugar."]),
            ("Três perfis", "Uma experiência objetiva para cada participante", ["Cliente: descoberta, reserva, cancelamento, avaliação, mensagens e notificações.", "Fornecedor: serviços, disponibilidade semanal, exceções, reservas e operação.", "Administração: status de fornecedores, KPIs, relatórios, disputas e categorias."]),
            ("Jornada do cliente", "Fluxo claro e auditável", ["Buscar e abrir o perfil de um fornecedor.", "Escolher serviço, data e horário disponível.", "Criar uma reserva temporária de cinco minutos durante o pagamento.", "Confirmar, receber notificações e avaliar após uma reserva elegível."]),
            ("Confiabilidade da reserva", "O banco é a autoridade final", ["A constraint PostgreSQL EXCLUDE USING gist bloqueia sobreposição de reservas ativas.", "A API traduz conflitos para uma resposta previsível ao cliente.", "Holds expirados são liberados pelo ciclo de vida do domínio.", "O modelo foi desenhado para permanecer correto sob concorrência."]),
            ("Arquitetura", "Camadas separadas e módulos testáveis", ["Next.js 14 para jornadas web por perfil.", "NestJS 11 para auth, catálogo, disponibilidade, reservas, pagamentos e admin.", "Prisma 5 e PostgreSQL 16 para transações, schema e migrations.", "Mock no desenvolvimento; fronteira Stripe para produção configurada."]),
            ("Base de segurança", "Proteção faz parte do contrato", ["Roles e ownership / IDOR checks para ações sensíveis.", "bcrypt, revogação de refresh token, throttling e validação na API.", "Cobertura de CSRF para mutações autenticadas por cookie.", "Raw body para webhook e transições idempotentes de pagamento."]),
            ("Evidência de qualidade", "O que foi verificado independentemente", ["16 de julho de 2026: 15 suites de API aprovadas.", "220 testes de API aprovados, sem snapshots com falha.", "A evidência demonstra comportamento testado, não deploy público.", "QA de navegador, providers externos e deploy permanecem gates separados."]),
            ("Limites da demo", "Transparência cria confiança", ["PAYMENTS_PROVIDER=mock é apenas para desenvolvimento e demonstração.", "Google (Demo) é uma simulação local; não é OAuth.", "Stripe real exige secrets, webhooks, deploy e testes E2E.", "O projeto não é anunciado como marketplace público com pagamentos reais."]),
            ("Próximos gates", "Do MVP entregue ao release candidate", ["Rodar aceitação de navegador em árabe RTL e inglês LTR.", "Validar Stripe test mode, e-mail e webhooks em ambiente adequado.", "Adicionar CI, deploy repetível, backup/restore, monitoring e performance.", "Coletar feedback antes de expandir o escopo comercial."]),
            ("Obrigado", "ServiceHub", ["Descoberta mais clara. Agendamento mais seguro. Operação mais inteligente.", "Documentação em inglês, árabe e português (Brasil)."]),
        ],
    },
}


def set_rtl(paragraph, enabled):
    if enabled:
        paragraph._p.get_or_add_pPr().set("rtl", "1")


def add_text(slide, x, y, w, h, value, size, color=TEXT, bold=False, align=PP_ALIGN.LEFT, rtl=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(.08); tf.margin_right = Inches(.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = value; p.alignment = align; set_rtl(p, rtl)
    for run in p.runs:
        run.font.name = FONT; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return shape


def rounded(slide, x, y, w, h, color=SURFACE, border=MINT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.color.rgb = border
    return shape


def build(language, deck):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    rtl = deck["rtl"]
    align = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT
    for number, (title, subtitle, bullets) in enumerate(deck["slides"], 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = BG
        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(.12))
        top.fill.solid(); top.fill.fore_color.rgb = MINT; top.line.fill.background()
        if number == 1:
            add_text(slide, 1.0, 1.25, 11.3, .8, title, 44, TEXT, True, PP_ALIGN.CENTER, rtl)
            add_text(slide, 1.0, 2.2, 11.3, .85, subtitle, 24, MINT, True, PP_ALIGN.CENTER, rtl)
            for i, bullet in enumerate(bullets):
                rounded(slide, 2.0 + i * 4.7, 4.5, 4.1, .8, SURFACE, GOLD if i else MINT)
                add_text(slide, 2.1 + i * 4.7, 4.6, 3.9, .55, bullet, 16, TEXT, True, PP_ALIGN.CENTER, rtl)
        elif number == len(deck["slides"]):
            add_text(slide, 1.0, 1.65, 11.3, .7, title, 42, MINT, True, PP_ALIGN.CENTER, rtl)
            add_text(slide, 1.0, 2.55, 11.3, .55, subtitle, 25, TEXT, True, PP_ALIGN.CENTER, rtl)
            rounded(slide, 1.5, 4.0, 10.3, 1.1, SURFACE, GOLD)
            add_text(slide, 1.8, 4.2, 9.7, .65, "\n".join(bullets), 16, MUTED, False, PP_ALIGN.CENTER, rtl)
        else:
            add_text(slide, .72, .44, 11.9, .5, title, 28, TEXT, True, align, rtl)
            add_text(slide, .78, 1.05, 11.75, .35, subtitle, 14, MUTED, False, align, rtl)
            for i, bullet in enumerate(bullets):
                y = 1.72 + i * 1.15
                accent = [MINT, GOLD, LILAC, MINT][i % 4]
                if rtl:
                    rounded(slide, 9.95, y, 2.35, .72, accent, accent)
                    add_text(slide, 1.05, y + .07, 8.55, .55, bullet, 18, TEXT, False, PP_ALIGN.RIGHT, rtl)
                else:
                    rounded(slide, 1.0, y, .72, .72, accent, accent)
                    add_text(slide, 2.0, y + .06, 10.1, .57, bullet, 18, TEXT, False, PP_ALIGN.LEFT, rtl)
            add_text(slide, .6, 7.02, 3.4, .22, f"{number:02d}  ·  ServiceHub", 10, MUTED, False, PP_ALIGN.LEFT, False)
            add_text(slide, 9.1, 7.02, 3.6, .22, deck["label"], 10, MUTED, False, PP_ALIGN.RIGHT, False)
    output = OUT_DIR / deck["file"]
    prs.save(output)
    print(output)


if __name__ == "__main__":
    for code, deck in DECKS.items():
        build(code, deck)
